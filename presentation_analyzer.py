from pptx import Presentation
from openai import OpenAI
import streamlit as st
import re
import PyPDF2

class PresentationAnalyzer:
    def __init__(self):
        self.client = OpenAI(api_key=st.secrets["openai_api_key"])
    
    def extract_text_from_ppt(self, file_path):
        """
        Estrae il testo da un file PowerPoint
        """
        prs = Presentation(file_path)
        text_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        return "\n".join(text_content)

    def extract_text_from_pdf(self, file_path):
        """
        Estrae il testo da un file PDF
        """
        text_content = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text_content.append(page.extract_text())
        
        return "\n".join(text_content)

    def extract_text(self, file_path):
        """
        Estrae il testo dal file in base all'estensione
        """
        file_extension = file_path.lower().split('.')[-1]
        
        if file_extension in ['ppt', 'pptx']:
            return self.extract_text_from_ppt(file_path)
        elif file_extension == 'pdf':
            return self.extract_text_from_pdf(file_path)
        else:
            raise ValueError(f"Formato file non supportato: {file_extension}")

    def analyze_with_gpt(self, text_content):
        """
        Analizza il testo estratto usando GPT per ottenere i campi richiesti
        """
        prompt = r"""Analizza questa presentazione PowerPoint e fornisci una risposta strutturata con questi elementi:

        1. CATEGORIA: [specifica una sola categoria tra: ACCESSO E POLICY MAKING, AWARENESS, EMPOWERMENT, PATIENT EXPERIENCE, PATIENT SUPPORT PROGRAM]
        
        2. INFORMAZIONI NECESSARIE ALLA GIURIA: [fornisci una descrizione molto dettagliata del progetto, includendo obiettivi e risultati.]
        
        3. SINTESI INFORMAZIONI PER L'EBOOK: [riassumi le informazioni principali del progetto in massimo 8 frasi, raccontando obiettivi e risultati brevemente]
        
        4. OBIETTIVI: [elenca tutti i principali obiettivi del progetto in modo conciso. Metti ogni obiettivo su una nuova riga, iniziando con - e andando a capo dopo ogni obiettivo]
        
        5. RISULTATI: [elenca tutti i principali risultati raggiunti in modo conciso. Metti ogni risultato su una nuova riga, iniziando con - e andando a capo dopo ogni risultato]

        Usa esattamente questi delimitatori nella tua risposta:
        <CATEGORIA>categoria</CATEGORIA>
        <INFO_GIURIA>informazioni complete</INFO_GIURIA>
        <SINTESI_EBOOK>sintesi breve</SINTESI_EBOOK>
        <OBIETTIVI>
        - primo obiettivo
        - secondo obiettivo
        (e così via per tutti gli obiettivi trovati)
        </OBIETTIVI>
        <RISULTATI>
        - primo risultato
        - secondo risultato
        (e così via per tutti i risultati trovati)
        </RISULTATI>

        IMPORTANTE: 
        - Non usare MAI virgolette doppie (") nel testo
        - Se devi citare qualcosa, usa solo virgolette singole (')
        - Non usare MAI il carattere backslash (\)
        - Non usare MAI <br> per andare a capo
        - Usa un vero ritorno a capo dopo ogni elemento con -
        - Elenca TUTTI gli obiettivi e risultati trovati
        
        Contenuto della presentazione:
        """ + text_content
        
        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Sei un assistente esperto nell'analisi di presentazioni in formato PowerPoint o PDF. Fornisci riassunti strutturati e completi."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=4096,
            temperature=0.7
        )
        
        return response.choices[0].message.content

    def extract_sections(self, gpt_response):
        """
        Estrae le sezioni dalla risposta di GPT
        """
        categoria = re.search(r'<CATEGORIA>(.*?)</CATEGORIA>', gpt_response, re.DOTALL)
        info_giuria = re.search(r'<INFO_GIURIA>(.*?)</INFO_GIURIA>', gpt_response, re.DOTALL)
        sintesi_ebook = re.search(r'<SINTESI_EBOOK>(.*?)</SINTESI_EBOOK>', gpt_response, re.DOTALL)
        obiettivi = re.search(r'<OBIETTIVI>(.*?)</OBIETTIVI>', gpt_response, re.DOTALL)
        risultati = re.search(r'<RISULTATI>(.*?)</RISULTATI>', gpt_response, re.DOTALL)
        
        return {
            "categoria": categoria.group(1).strip() if categoria else "Non specificata",
            "info_giuria": info_giuria.group(1).strip() if info_giuria else "Non disponibile",
            "sintesi_ebook": sintesi_ebook.group(1).strip() if sintesi_ebook else "Non disponibile",
            "obiettivi": obiettivi.group(1).strip() if obiettivi else "Non specificati",
            "risultati": risultati.group(1).strip() if risultati else "Non specificati"
        }

    def analyze(self, file_path):
        """
        Funzione principale che coordina l'estrazione e l'analisi della presentazione
        """
        try:
            text_content = self.extract_text(file_path)
            gpt_response = self.analyze_with_gpt(text_content)
            return self.extract_sections(gpt_response)
        except Exception as e:
            st.error(f"Errore durante l'analisi della presentazione: {e}")
            return None